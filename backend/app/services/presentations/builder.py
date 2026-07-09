"""Builds a .pptx from an AI-generated PresentationOutline using python-pptx's built-in default
template. There's no real company branding (logo/colors) to work from yet, so this applies a
single accent color rather than a bespoke branded template — swap PRIMARY_COLOR (and add a real
template file) once brand assets exist.
"""

import io

from pptx import Presentation as PptxPresentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from app.services.ai.schemas import PresentationOutline, SlideChart, SlideContent

PRIMARY_COLOR = RGBColor(0x1E, 0x29, 0x3B)  # slate navy, matches the web dashboard's primary color

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

_TITLE_LAYOUT = 0
_CONTENT_LAYOUT = 1
_SECTION_LAYOUT = 2
_TWO_CONTENT_LAYOUT = 3
_BLANK_LAYOUT = 6


def build_presentation(outline: PresentationOutline) -> bytes:
    prs = PptxPresentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _add_title_slide(prs, outline.title)
    for slide_content in outline.slides:
        _add_content_slide(prs, slide_content)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _add_title_slide(prs: PptxPresentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_LAYOUT])
    _style_title(slide.shapes.title, title, size=Pt(40))
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text_frame.text = "Prepared by AI Executive Assistant"


def _add_content_slide(prs: PptxPresentation, slide_content: SlideContent) -> None:
    if slide_content.layout == "section_header":
        slide = prs.slides.add_slide(prs.slide_layouts[_SECTION_LAYOUT])
        _style_title(slide.shapes.title, slide_content.title)
    elif slide_content.layout == "two_column":
        slide = prs.slides.add_slide(prs.slide_layouts[_TWO_CONTENT_LAYOUT])
        _style_title(slide.shapes.title, slide_content.title)
        _fill_bullets(slide.placeholders[1], slide_content.left_column)
        _fill_bullets(slide.placeholders[2], slide_content.right_column)
    elif slide_content.layout == "chart" and slide_content.chart:
        slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
        _add_title_textbox(slide, slide_content.title)
        _add_chart(slide, slide_content.chart)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[_CONTENT_LAYOUT])
        _style_title(slide.shapes.title, slide_content.title)
        _fill_bullets(slide.placeholders[1], slide_content.bullets)

    if slide_content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_content.speaker_notes


def _style_title(title_shape, text: str, size: Pt | None = None) -> None:
    if title_shape is None:
        return
    title_shape.text_frame.text = text
    run = title_shape.text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    run.font.color.rgb = PRIMARY_COLOR
    if size is not None:
        run.font.size = size


def _add_title_textbox(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), SLIDE_WIDTH - Inches(1), Inches(1))
    text_frame = box.text_frame
    text_frame.text = text
    run = text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = PRIMARY_COLOR


def _fill_bullets(placeholder, bullets: list[str]) -> None:
    text_frame = placeholder.text_frame
    text_frame.clear()
    items = bullets or ["(no content provided)"]
    text_frame.text = items[0]
    for bullet in items[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = bullet


def _add_chart(slide, chart: SlideChart) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = chart.categories
    chart_data.add_series(chart.series_name, chart.values)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1.5),
        SLIDE_WIDTH - Inches(2),
        SLIDE_HEIGHT - Inches(2.5),
        chart_data,
    )
